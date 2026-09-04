from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1F, 0x38, 0x64)
STEEL     = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLU = RGBColor(0xDF, 0xE9, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GOLD      = RGBColor(0xF0, 0xA5, 0x00)
DARK      = RGBColor(0x1A, 0x1A, 0x2E)
DGRAY     = RGBColor(0x33, 0x33, 0x33)
MGRAY     = RGBColor(0x55, 0x55, 0x55)
LGRAY     = RGBColor(0xF4, 0xF6, 0xFB)
GREEN     = RGBColor(0x1B, 0x5E, 0x20)
LGREEN    = RGBColor(0xE8, 0xF5, 0xE9)
ORANGE    = RGBColor(0xBF, 0x36, 0x0C)
LORANGE   = RGBColor(0xFF, 0xF3, 0xE0)
PURPLE    = RGBColor(0x4A, 0x14, 0x8C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

W = 13.33   # slide width
H = 7.5     # slide height

# ── Low-level helpers ────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    from pptx.util import Emu
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.width = lw
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, sz=16, bold=False, color=DGRAY,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return box

def mtb(slide, lines, l, t, w, h, default_sz=15,
        default_bold=False, default_color=DGRAY, default_align=PP_ALIGN.LEFT):
    """lines = list of str  OR  (text, bold, color, size, align)"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, s, a = item, default_bold, default_color, default_sz, default_align
        else:
            txt = item[0]
            b   = item[1] if len(item) > 1 else default_bold
            c   = item[2] if len(item) > 2 else default_color
            s   = item[3] if len(item) > 3 else default_sz
            a   = item[4] if len(item) > 4 else default_align
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = a
        run = p.add_run(); run.text = txt
        run.font.size = Pt(s); run.font.bold = b; run.font.color.rgb = c

def bullets(slide, items, l, t, w, h, sz=15):
    lines = []
    for item in items:
        sub = item.startswith("    ")
        txt = item.strip()
        bul = "  -" if sub else "•"
        c   = MGRAY if sub else DGRAY
        s   = sz - 1 if sub else sz
        lines.append((f"{bul}  {txt}", False, c, s, PP_ALIGN.LEFT))
    mtb(slide, lines, l, t, w, h, default_sz=sz)

def tbl(slide, headers, rows, l, t, w, h, hbg=STEEL, hfg=WHITE,
        alt=LIGHT_BLU, fs=13):
    cols  = len(headers)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, cols,
                Inches(l), Inches(t), Inches(w), Inches(h))
    table = shape.table
    col_w = Inches(w / cols)
    for c in range(cols):
        table.columns[c].width = col_w

    def fmt(cell, text, bg, fg, bold=False):
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.size = Pt(fs); run.font.bold = bold
        run.font.color.rgb = fg

    for c, h in enumerate(headers):
        fmt(table.cell(0, c), h, hbg, hfg, bold=True)
    for r, row in enumerate(rows):
        bg = WHITE if r % 2 == 0 else alt
        for c, val in enumerate(row):
            fmt(table.cell(r+1, c), val, bg, DGRAY)

# ── Reusable slide chrome ────────────────────────────────────────────────────
def bg(slide):
    rect(slide, 0, 0, W, H, fill=LGRAY)

def header(slide, title, subtitle=None):
    """Dark navy top bar."""
    rect(slide, 0, 0, W, 1.30, fill=NAVY)
    rect(slide, 0, 1.30, W, 0.055, fill=GOLD)
    tb(slide, title, 0.40, 0.10, W-0.8, 0.80,
       sz=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        tb(slide, subtitle, 0.40, 0.82, W-0.8, 0.44,
           sz=14, color=GOLD, align=PP_ALIGN.LEFT)

def section_pill(slide, label, l, t, w=3.5, h=0.38,
                 bg_col=STEEL, fg_col=WHITE):
    rect(slide, l, t, w, h, fill=bg_col)
    tb(slide, label, l, t, w, h, sz=13, bold=True,
       color=fg_col, align=PP_ALIGN.CENTER)

def info_box(slide, title, body_lines, l, t, w, h,
             hdr_col=STEEL, body_col=LGRAY, sz=13):
    rect(slide, l, t, w, h, fill=body_col, line=hdr_col, lw=Pt(1))
    rect(slide, l, t, w, 0.38, fill=hdr_col)
    tb(slide, title, l, t, w, 0.38, sz=13, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    bullets(slide, body_lines, l+0.1, t+0.40, w-0.2, h-0.45, sz=sz)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=DARK)
rect(s, 0, 0, W, H*0.55, fill=NAVY)   # top half
rect(s, 0, H*0.55-0.04, W, 0.08, fill=GOLD)

tb(s, "COMSATS University Islamabad",
   0.6, 0.22, W-1.2, 0.5, sz=14, color=LIGHT_BLU,
   align=PP_ALIGN.CENTER)

tb(s, "Dynamic Target Assignment and\nCollision Avoidance for Multi-UAV\nCoordination in 3D using MARL",
   0.6, 0.72, W-1.2, 2.6, sz=30, bold=True,
   color=WHITE, align=PP_ALIGN.CENTER)

tb(s, "MS Research Synopsis Presentation",
   0.6, 3.28, W-1.2, 0.5, sz=16, color=GOLD, align=PP_ALIGN.CENTER)

# bottom strip
mtb(s, [
    ("Ayesha Khalil", True, WHITE, 20, PP_ALIGN.CENTER),
    ("SP25-RCS-009   |   MS Computer Science", False, LIGHT_BLU, 14, PP_ALIGN.CENTER),
    ("Department of Computer Science", False, LIGHT_BLU, 13, PP_ALIGN.CENTER),
], 0.6, 4.35, W-1.2, 1.2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – OUTLINE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Presentation Outline")

sections = [
    (STEEL,  "01", "Introduction",          "Motivation, real-world context, and why multi-UAV coordination matters"),
    (STEEL,  "02", "Background",            "Deep Reinforcement Learning, MARL, and the CTDE framework"),
    (GREEN,  "03", "Literature Review",     "10 reviewed papers grouped by contribution area"),
    (GREEN,  "04", "Most Related Work",     "DA-MAPPO and IGAT-MARL in detail"),
    (ORANGE, "05", "Problem Statement",     "The identified gap and the interference problem in 3D"),
    (ORANGE, "06", "Research Objectives",   "Four measurable research objectives"),
    (PURPLE, "07", "Proposed Methodology",  "Framework design, observation vector, training strategy"),
    (PURPLE, "08", "Evaluation Plan",       "Metrics, baselines, and test configurations"),
    (NAVY,   "09", "Expected Contributions","What this research will add to the field"),
]
col_gap = 0.25
box_w   = (W - 0.5 - col_gap) / 2
for i, (col, num, title, desc) in enumerate(sections):
    col_idx = i % 2
    row_idx = i // 2
    x = 0.25 + col_idx * (box_w + col_gap)
    y = 1.50 + row_idx * 1.35
    if i == 8:   # last item – center it
        x = (W - box_w) / 2
    rect(s, x, y, box_w, 1.15, fill=WHITE, line=col, lw=Pt(1.5))
    rect(s, x, y, 0.55, 1.15, fill=col)
    tb(s, num, x, y, 0.55, 1.15, sz=16, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, x+0.62, y+0.06, box_w-0.7, 0.42,
       sz=14, bold=True, color=col)
    tb(s, desc,  x+0.62, y+0.48, box_w-0.7, 0.60,
       sz=11, color=MGRAY, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – INTRODUCTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Introduction",
       "Unmanned Aerial Vehicles (UAVs) in time-critical missions")

bullets(s, [
    "UAVs have become essential tools in disaster response, search and rescue, infrastructure inspection, and surveillance",
    "These applications share a common requirement: not one drone working alone, but teams of drones coordinating in real-time",
    "Each drone in such a team must simultaneously navigate toward an assigned target, avoid collisions with other drones, and share limited airspace",
    "Classical path-planning algorithms such as A* and RRT require a fully known, static environment and cannot adapt when conditions change mid-flight",
    "Deep Reinforcement Learning (DRL) learns navigation policies through interaction with a simulated environment, enabling real-time adaptation to dynamic conditions",
    "Multi-Agent Deep Reinforcement Learning (MARL) extends this to teams of agents, where each drone learns to act cooperatively with others",
    "Despite rapid progress, two critical sub-problems have been developed in isolation from each other:",
    "    Dynamic target assignment — which drone chases which target, updated in real-time",
    "    Collision avoidance — preventing mid-air collisions between coordinating drones",
    "No published work has integrated both mechanisms into a single policy, and none has been tested in a 3D environment",
], 0.40, 1.50, 8.3, 5.7, sz=15)

# side panel
rect(s, 9.0, 1.50, 4.0, 5.7, fill=WHITE, line=STEEL, lw=Pt(1.2))
rect(s, 9.0, 1.50, 4.0, 0.42, fill=NAVY)
tb(s, "This Research Addresses", 9.0, 1.50, 4.0, 0.42,
   sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

points = [
    ("Real-time target assignment", STEEL),
    ("Collision avoidance in 3D", GREEN),
    ("Both in a single policy", ORANGE),
    ("Up to 8 drones simultaneously", PURPLE),
    ("First empirical results in 3D", NAVY),
]
for i, (pt, c) in enumerate(points):
    rect(s, 9.1, 2.05+i*0.95, 3.8, 0.75, fill=LIGHT_BLU, line=c, lw=Pt(1.2))
    tb(s, pt, 9.1, 2.05+i*0.95, 3.8, 0.75,
       sz=13, color=c, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – BACKGROUND: DRL
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Background: Deep Reinforcement Learning",
       "The foundation of adaptive UAV navigation")

# core DRL loop diagram (3 boxes with arrows)
boxes = [
    (0.40, 3.0,  "Agent\n(Neural Network Policy)", NAVY),
    (4.65, 3.0,  "Environment\n(3D Simulation)", STEEL),
    (4.65, 5.2,  "Reward Signal\n(r)", GOLD),
]
bw, bh = 3.5, 1.2
for (bx, by, label, col) in boxes:
    rect(s, bx, by, bw, bh, fill=col)
    tb(s, label, bx, by, bw, bh, sz=15, bold=True,
       color=WHITE if col != GOLD else DARK, align=PP_ALIGN.CENTER)

# arrows
tb(s, "→  Action (a)", 3.95, 3.32, 0.80, 0.4, sz=11, color=DGRAY)
tb(s, "←  State (s)",  3.95, 3.72, 0.80, 0.4, sz=11, color=DGRAY)
tb(s, "↓",  6.15, 4.20, 0.5, 0.5, sz=18, bold=True, color=STEEL, align=PP_ALIGN.CENTER)
tb(s, "↑  Update policy", 2.0, 5.55, 2.2, 0.4, sz=11, color=DGRAY)

# key concepts on the right
info_box(s, "Key DRL Concepts",
    ["Policy: the neural network that maps state to action",
     "Reward: numerical feedback signal (+/- per step)",
     "Value function: expected future reward from a state",
     "Episode: one complete mission run from start to end",
     "Curriculum learning: train on easy cases first, then harder ones",
    ], 8.6, 1.55, 4.45, 3.0, sz=13)

info_box(s, "Why DRL for UAVs?",
    ["Handles dynamic, partially observable environments",
     "No need for hand-coded rules — learns from experience",
     "Policy improves over millions of simulation episodes",
     "Generalizes to situations not seen during training",
    ], 8.6, 4.70, 4.45, 2.50, sz=13)

# left side concepts
info_box(s, "CTDE Framework",
    ["Centralized Training: critic sees all agents' states during training",
     "Decentralized Execution: each drone acts on its own observations only",
     "Solves the non-stationarity problem in multi-agent learning",
    ], 0.40, 1.55, 7.9, 1.55, sz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – BACKGROUND: MARL
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Background: Multi-Agent RL and Key Algorithms",
       "Algorithms directly relevant to this research")

section_pill(s, "Multi-Agent PPO (MAPPO)", 0.40, 1.50, 4.5)
bullets(s, [
    "PPO extended to multi-agent cooperative settings",
    "Shared centralized critic trained on global state",
    "Each agent executes using only its local observation",
    "Stable, reproducible, competitive with more complex MARL methods",
    "Used as the policy backbone in this research",
], 0.40, 1.93, 4.5, 2.5, sz=14)

section_pill(s, "Hungarian Algorithm", 5.10, 1.50, 4.0)
bullets(s, [
    "Optimal one-to-one assignment algorithm",
    "Assigns N drones to N targets at minimum total cost",
    "Runs at every decision step in this research",
    "Result added directly to each drone's observation vector",
    "Computational complexity: O(N^3) — feasible for small swarms",
], 5.10, 1.93, 4.0, 2.5, sz=14)

section_pill(s, "Conflict Graph (Sparse Interaction)", 9.30, 1.50, 3.75)
bullets(s, [
    "Connects only drone pairs predicted to collide",
    "Much fewer edges than all-to-all communication",
    "Updated dynamically every decision step",
    "Used for collision avoidance in this research",
], 9.30, 1.93, 3.75, 2.5, sz=14)

rect(s, 0.40, 4.50, W-0.8, 0.05, fill=GOLD)

section_pill(s, "Dec-POMDP: The Formal Framework", 0.40, 4.62, 5.5)
bullets(s, [
    "Decentralized Partially Observable Markov Decision Process",
    "Each drone has its own local observation — cannot see everything",
    "Agents cooperate to maximize a shared team reward",
    "Standard mathematical formulation for all MARL UAV problems",
], 0.40, 5.05, 5.5, 2.2, sz=14)

section_pill(s, "MAPPO vs MADDPG vs QMIX", 6.15, 4.62, 6.8)
tb(s, "This research uses MAPPO because:", 6.15, 5.07, 6.8, 0.35,
   sz=13, bold=True, color=DGRAY)
tbl(s,
    ["Algorithm", "Type", "Chosen?"],
    [["MAPPO",   "On-policy, stable, continuous/discrete", "YES — used here"],
     ["MADDPG",  "Off-policy, continuous only", "No"],
     ["QMIX",    "Value decomposition, discrete actions", "No"],
    ],
    6.15, 5.42, 6.8, 1.88, fs=12)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – LITERATURE REVIEW: Single-UAV Navigation
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Literature Review",
       "10 papers reviewed — grouped by contribution area")

# section label
rect(s, 0.40, 1.50, W-0.8, 0.42, fill=NAVY)
tb(s, "GROUP 1: Single-UAV Navigation   [Papers 1 & 3]",
   0.40, 1.50, W-0.8, 0.42, sz=15, bold=True, color=WHITE)

# Paper 1
rect(s, 0.40, 2.05, 6.1, 2.55, fill=WHITE, line=STEEL, lw=Pt(1.2))
rect(s, 0.40, 2.05, 6.1, 0.40, fill=STEEL)
tb(s, "[1]  Tang, Liang & Li (2024) — Drones Journal",
   0.50, 2.05, 6.0, 0.40, sz=12, bold=True, color=WHITE)
mtb(s, [
    ("Method: ", True, NAVY, 13), ("Improved D3QN = Dueling Double DQN + Prioritized Experience Replay + heuristic bias toward target", False, DGRAY, 13),
], 0.50, 2.52, 5.9, 0.4)
mtb(s, [
    ("Setting: ", True, NAVY, 13), ("Single UAV, 2D environment, static and moving obstacles", False, DGRAY, 13),
], 0.50, 2.88, 5.9, 0.4)
mtb(s, [
    ("Result: ", True, GREEN, 13), ("95% success rate — outperforms A* and RRT on path quality", False, DGRAY, 13),
], 0.50, 3.24, 5.9, 0.4)
mtb(s, [
    ("Limitation: ", True, ORANGE, 13), ("Single drone only, 2D only, no multi-agent extension", False, DGRAY, 13),
], 0.50, 3.60, 5.9, 0.4)

# Paper 3
rect(s, 6.80, 2.05, 6.15, 2.55, fill=WHITE, line=STEEL, lw=Pt(1.2))
rect(s, 6.80, 2.05, 6.15, 0.40, fill=STEEL)
tb(s, "[3]  Jarray, Zaghbani & Bouallègue (2025) — Procedia Computer Science",
   6.90, 2.05, 6.05, 0.40, sz=12, bold=True, color=WHITE)
mtb(s, [
    ("Method: ", True, NAVY, 13), ("DQN + 3D CNN + dynamic reward formula d3/(d1+d2) per decision step", False, DGRAY, 13),
], 6.90, 2.52, 5.95, 0.4)
mtb(s, [
    ("Setting: ", True, NAVY, 13), ("Single UAV, large-scale 3D environment (25 km²), 4 obstacle densities", False, DGRAY, 13),
], 6.90, 2.88, 5.95, 0.4)
mtb(s, [
    ("Result: ", True, GREEN, 13), ("98% success in low obstacles, 85% in high obstacles — beats PSO and GWO", False, DGRAY, 13),
], 6.90, 3.24, 5.95, 0.4)
mtb(s, [
    ("Limitation: ", True, ORANGE, 13), ("Only one drone — the only paper with true 3D but cannot coordinate teams", False, DGRAY, 13),
], 6.90, 3.60, 5.95, 0.4)

rect(s, 0.40, 4.72, W-0.8, 0.42, fill=LIGHT_BLU, line=STEEL, lw=Pt(1))
tb(s, "Gap identified:  Both papers solve navigation for a SINGLE drone. Neither scales to multiple drones or addresses coordination.",
   0.55, 4.72, W-1.1, 0.42, sz=13, bold=False, color=NAVY)

# GROUP 2 label
rect(s, 0.40, 5.22, W-0.8, 0.42, fill=NAVY)
tb(s, "GROUP 2: Multi-UAV Target Assignment   [Papers 2 & 4]",
   0.40, 5.22, W-0.8, 0.42, sz=15, bold=True, color=WHITE)

rect(s, 0.40, 5.72, 6.1, 1.55, fill=WHITE, line=STEEL, lw=Pt(1))
tb(s, "[2]  Kong et al. (2024) — Frontiers in Neurorobotics",
   0.50, 5.72, 6.0, 0.35, sz=12, bold=True, color=NAVY)
bullets(s, ["TANet-TD3: Target Assignment Network + TD3 policy",
            "5 UAVs, 2D, partial observability, Hungarian supervision",
            "Limitation: 5 drones only, 2D, no scaling mechanism"],
        0.50, 6.10, 5.9, 1.1, sz=13)

rect(s, 6.80, 5.72, 6.15, 1.55, fill=WHITE, line=STEEL, lw=Pt(1))
tb(s, "[4]  Zhang et al. (2025) — Chinese Journal of Aeronautics",
   6.90, 5.72, 6.05, 0.35, sz=12, bold=True, color=NAVY)
bullets(s, ["PO-WMFDDPG: Mean-field + multi-head attention + CTDE",
            "Scales to 120 drones, >90% success — 2D only",
            "Limitation: 2D, static obstacles, no dynamic reassignment"],
        6.90, 6.10, 5.95, 1.1, sz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – LITERATURE REVIEW: LLM-MARL + Survey
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Literature Review (Continued)",
       "LLM-enhanced coordination and survey literature")

rect(s, 0.40, 1.50, W-0.8, 0.42, fill=NAVY)
tb(s, "GROUP 3: LLM-Enhanced Multi-Agent Coordination   [Papers 6, 7 & 8]",
   0.40, 1.50, W-0.8, 0.42, sz=15, bold=True, color=WHITE)

paper_data = [
    ("[6]  Xu et al. (2026) — arXiv 2505.08448",
     "MRLMN: IPPO + role-based reward + GPT-4o knowledge distillation via Hungarian matching",
     "12-24 UAVs, 52% higher data rate, 27% more user coverage",
     "Simulation only, GPT-4o dependency, 2D movement only"),
    ("[7]  Wang et al. (2025) — IEEE Open Journal of Vehicular Technology",
     "RALLY: Two-stage — LLM semantic consensus assigns Commander/Coordinator/Executor roles, then RMIX value-mixing optimizes coordination",
     "Generalizes to unseen swarm sizes without retraining; < 5 GB deployment",
     "14.48s inference latency, qualitative validation only, no real hardware"),
    ("[8]  Khan et al. (2026) — IEEE Access",
     "Systematic comparison of CrewAI, AutoGen, LangChain + LangGraph-CrewAI hybrid with complexity-aware routing",
     "96.1% success at 76% lower token cost; first fair large-scale framework comparison",
     "Wildfire domain only; hybrid validated with only 6 agents"),
]
for i, (title, method, result, limit) in enumerate(paper_data):
    x = 0.40 + i * 4.30
    rect(s, x, 2.05, 4.1, 2.60, fill=WHITE, line=STEEL, lw=Pt(1))
    rect(s, x, 2.05, 4.1, 0.38, fill=STEEL)
    tb(s, title, x+0.08, 2.05, 3.95, 0.38, sz=11, bold=True, color=WHITE)
    mtb(s, [("Method: ", True, NAVY, 12), (method, False, DGRAY, 12)],
        x+0.08, 2.48, 3.94, 0.7)
    mtb(s, [("Result: ", True, GREEN, 12), (result, False, DGRAY, 12)],
        x+0.08, 3.15, 3.94, 0.6)
    mtb(s, [("Limit: ", True, ORANGE, 12), (limit, False, DGRAY, 12)],
        x+0.08, 3.72, 3.94, 0.6)

rect(s, 0.40, 4.78, W-0.8, 0.40, fill=LIGHT_BLU, line=STEEL, lw=Pt(1))
tb(s, "Gap identified:  LLM approaches improve coordination strategy but none study the specific interference between target assignment and collision avoidance in 3D.",
   0.55, 4.78, W-1.1, 0.40, sz=13, color=NAVY)

rect(s, 0.40, 5.28, W-0.8, 0.40, fill=NAVY)
tb(s, "GROUP 4: Survey Literature   [Paper 11]",
   0.40, 5.28, W-0.8, 0.40, sz=15, bold=True, color=WHITE)

rect(s, 0.40, 5.75, W-0.8, 1.58, fill=WHITE, line=STEEL, lw=Pt(1))
tb(s, "[11]  Govinda, Brik & Harous (2025) — IEEE Transactions on Intelligent Transportation Systems",
   0.55, 5.75, W-1.1, 0.38, sz=13, bold=True, color=NAVY)
bullets(s, [
    "Systematic literature review of DRL applications across transportation, robotics, and UAV systems",
    "Key finding: lack of unified frameworks combining navigation efficiency and inter-agent coordination is a prevailing open gap",
    "Directly motivates this research — the gap they identified is exactly what this work addresses",
], 0.55, 6.17, W-1.1, 1.1, sz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – MOST RELATED WORK: DA-MAPPO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Most Related Work: DA-MAPPO  [Paper 10]",
       "Sheng, Xie, Liu & Li (2026) — IEEE Internet of Things Journal")

# left: what it does
rect(s, 0.40, 1.52, 6.1, 5.75, fill=WHITE, line=STEEL, lw=Pt(1.2))
section_pill(s, "What DA-MAPPO Does", 0.40, 1.52, 6.1, 0.40, STEEL)
bullets(s, [
    "Solves real-time dynamic target assignment for UAV swarms",
    "Hungarian algorithm runs at every single decision step",
    "Assignment result encoded directly into each drone's observation",
    "Policy: MAPPO with 4-tier hierarchical reward structure",
    "Trained with curriculum: 3 difficulty levels, 3 to 5 drones",
    "Achieves 90-99% mission success — 25 pp above best baseline",
    "Tested with 3 UAVs, 3 dynamic targets, 30-50 static obstacles",
    "Validated on NVIDIA Jetson Orin Nano embedded hardware",
], 0.40, 1.95, 6.1, 4.0, sz=14)

section_pill(s, "What DA-MAPPO Leaves Unresolved", 0.40, 6.05, 6.1, 0.40, ORANGE)
bullets(s, [
    "Only 2D — no altitude, no vertical flight paths",
    "Only 3 drones — scalability unknown",
    "No explicit collision-avoidance mechanism (only a reward penalty)",
    "Authors explicitly list 3D extension as future work",
], 0.40, 6.48, 6.1, 1.0, sz=13)

# right: critical finding box
rect(s, 6.80, 1.52, 6.15, 3.5, fill=RGBColor(0xFF,0xF8,0xE8),
     line=GOLD, lw=Pt(2))
rect(s, 6.80, 1.52, 6.15, 0.40, fill=ORANGE)
tb(s, "Critical Ablation Result", 6.80, 1.52, 6.15, 0.40,
   sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "With assignment mechanism:", 6.95, 2.00, 5.9, 0.35, sz=13, bold=True, color=GREEN)
rect(s, 6.95, 2.38, 5.9, 0.55, fill=GREEN)
tb(s, "Mission success: 90 - 99%", 6.95, 2.38, 5.9, 0.55,
   sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Without assignment mechanism:", 6.95, 3.05, 5.9, 0.35, sz=13, bold=True, color=ORANGE)
rect(s, 6.95, 3.43, 5.9, 0.55, fill=ORANGE)
tb(s, "Mission success: 0%", 6.95, 3.43, 5.9, 0.55,
   sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "The assignment mechanism is not just helpful —\nit is the entire reason the policy works.",
   6.95, 4.10, 5.9, 0.78, sz=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# right: architecture
rect(s, 6.80, 5.12, 6.15, 2.15, fill=WHITE, line=STEEL, lw=Pt(1))
section_pill(s, "Observation Vector in DA-MAPPO", 6.80, 5.12, 6.15, 0.38, STEEL)
bullets(s, [
    "Own position and velocity (3D position in 2D env.)",
    "Relative position to ASSIGNED target (Hungarian output)",
    "Positions of visible obstacles",
    "No conflict-neighbor information — no collision awareness",
], 6.80, 5.54, 6.15, 1.65, sz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – MOST RELATED WORK: IGAT-MARL
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Most Related Work: IGAT-MARL  [Paper 9]",
       "Rezaee, Abdul Hamid, Hussin & Zukarnain (2026) — Applied Soft Computing")

rect(s, 0.40, 1.52, 6.1, 5.75, fill=WHITE, line=STEEL, lw=Pt(1.2))
section_pill(s, "What IGAT-MARL Does", 0.40, 1.52, 6.1, 0.40, STEEL)
bullets(s, [
    "Solves collision avoidance between multiple UAVs at scale",
    "Builds a sparse conflict graph at every decision step",
    "Only connects drone pairs predicted to collide within a time window",
    "Replaces costly all-to-all communication with targeted interaction",
    "Uses Improved Graph Attention Network: stacked double attention + residual connections",
    "Trained with curriculum learning on 3 to 10 drones",
    "Simulation: BlueSky air traffic simulator (realistic fixed-wing dynamics)",
], 0.40, 1.95, 6.1, 3.6, sz=14)

section_pill(s, "What IGAT-MARL Leaves Unresolved", 0.40, 5.62, 6.1, 0.40, ORANGE)
bullets(s, [
    "No target assignment — drones have no goal to fly toward",
    "Fixed-wing aircraft only (not rotary UAVs)",
    "Simulation only — no hardware validation",
    "Authors explicitly list task allocation as future direction",
], 0.40, 6.05, 6.1, 1.2, sz=13)

# results
rect(s, 6.80, 1.52, 6.15, 3.0, fill=LGREEN, line=GREEN, lw=Pt(1.5))
rect(s, 6.80, 1.52, 6.15, 0.40, fill=GREEN)
tb(s, "Key Results vs. Best Baseline", 6.80, 1.52, 6.15, 0.40,
   sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
results = [
    ("17%", "higher cumulative reward"),
    ("10%", "fewer dangerous separation events"),
    ("44%", "fewer interaction graph edges"),
]
for i, (pct, label) in enumerate(results):
    rect(s, 7.00, 2.05+i*0.72, 5.75, 0.60, fill=GREEN)
    tb(s, f"{pct}  {label}", 7.00, 2.05+i*0.72, 5.75, 0.60,
       sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# sparse graph explanation
rect(s, 6.80, 4.65, 6.15, 2.62, fill=WHITE, line=STEEL, lw=Pt(1))
section_pill(s, "Why the Sparse Conflict Graph Matters", 6.80, 4.65, 6.15, 0.38, STEEL)
bullets(s, [
    "All-to-all graph: N drones = N*N edges — grows quadratically",
    "Conflict graph: only predicted collision pairs connected",
    "44% fewer edges means 44% less unnecessary computation and communication",
    "Critical for scaling to larger drone swarms",
    "This mechanism is what this research combines with DA-MAPPO",
], 6.80, 5.07, 6.15, 2.10, sz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Problem Statement",
       "The gap that exists between DA-MAPPO and IGAT-MARL")

# gap summary table
tbl(s,
    ["", "DA-MAPPO [10]", "IGAT-MARL [9]", "This Research"],
    [
        ["Target Assignment",    "YES (Hungarian, per step)", "NO",                "YES"],
        ["Collision Avoidance",  "Reward penalty only",       "YES (sparse graph)", "YES"],
        ["3D Environment",       "NO (2D only)",              "NO (2D only)",       "YES"],
        ["Swarm Size Tested",    "3 drones",                  "3-10 drones",        "3, 5, 8 drones"],
        ["Other cites as future","3D + larger swarms",        "Task allocation",    "Both combined"],
    ],
    0.40, 1.52, W-0.8, 2.90, fs=13)

# interference explanation
rect(s, 0.40, 4.55, W-0.8, 0.40, fill=NAVY)
tb(s, "The Interference Problem in 3D", 0.40, 4.55, W-0.8, 0.40,
   sz=15, bold=True, color=WHITE)

rect(s, 0.40, 5.00, 6.0, 2.32, fill=LORANGE, line=ORANGE, lw=Pt(1.5))
tb(s, "Assignment module says:", 0.55, 5.03, 5.8, 0.35,
   sz=13, bold=True, color=ORANGE)
tb(s, '"Fly directly toward your assigned target now"',
   0.55, 5.42, 5.8, 0.38, sz=13, italic=True, color=DARK)
tb(s, "But has no knowledge of which other drone pairs are in a conflict state",
   0.55, 5.84, 5.8, 0.55, sz=13, color=DGRAY)
tb(s, "In 3D: a drone assigned a target directly above it is on a vertical collision course with a descending drone — assignment module does not see this",
   0.55, 6.42, 5.8, 0.8, sz=12, color=ORANGE, bold=True)

rect(s, 6.80, 5.00, 6.15, 2.32, fill=LIGHT_BLU, line=STEEL, lw=Pt(1.5))
tb(s, "Conflict graph says:", 6.95, 5.03, 5.9, 0.35,
   sz=13, bold=True, color=STEEL)
tb(s, '"Change your course — collision predicted"',
   6.95, 5.42, 5.9, 0.38, sz=13, italic=True, color=DARK)
tb(s, "But has no knowledge of what the current target assignment is",
   6.95, 5.84, 5.9, 0.55, sz=13, color=DGRAY)
tb(s, "Rerouting causes course deviations that may pull drones away from their assigned targets, creating oscillation or mission failure",
   6.95, 6.42, 5.9, 0.8, sz=12, color=STEEL, bold=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – RESEARCH QUESTION & OBJECTIVES
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Research Question & Objectives")

rect(s, 0.40, 1.52, W-0.8, 1.50, fill=NAVY)
tb(s,
   "Will a single MAPPO policy trained with a conflict-aware interaction graph and assignment-augmented observation successfully navigate 5-8 UAVs to dynamically assigned missions in a 3D environment with mission success above 85%, and do the two mechanisms help or hinder each other?",
   0.60, 1.56, W-1.2, 1.42, sz=15, bold=True, color=GOLD,
   align=PP_ALIGN.CENTER)

rect(s, 0.40, 3.12, W-0.8, 0.40, fill=STEEL)
tb(s, "Four Research Objectives", 0.40, 3.12, W-0.8, 0.40,
   sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

objs = [
    (STEEL,  "Obj 1",
     "Design a combined observation vector that encodes both Hungarian assignment state and conflict-aware neighborhood for each drone inside a single 3D MAPPO policy."),
    (GREEN,  "Obj 2",
     "Evaluate the framework with 3, 5, and 8 drones to quantify how performance scales with increasing simultaneous assignment-conflict interactions."),
    (ORANGE, "Obj 3",
     "Run controlled ablation experiments isolating the contribution of the conflict graph, the assignment mechanism, and the 3D extension independently and in combination."),
    (PURPLE, "Obj 4",
     "Map the failure boundary: the region of swarm size, obstacle density, and target speed where the combined policy fails, and characterize the mode of failure."),
]
for i, (col, label, desc) in enumerate(objs):
    y = 3.65 + i * 0.92
    rect(s, 0.40, y, 1.35, 0.78, fill=col)
    tb(s, label, 0.40, y, 1.35, 0.78, sz=13, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)
    rect(s, 1.80, y, W-2.2, 0.78, fill=WHITE, line=col, lw=Pt(1.2))
    tb(s, desc, 1.95, y+0.08, W-2.4, 0.65, sz=13, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – PROPOSED METHODOLOGY
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Proposed Methodology: Framework Design",
       "A unified observation encoding assignment state and conflict neighborhood")

# Observation vector (left)
rect(s, 0.40, 1.52, 5.8, 5.75, fill=WHITE, line=STEEL, lw=Pt(1.2))
section_pill(s, "Combined Observation Vector per Drone", 0.40, 1.52, 5.8, 0.40, NAVY)
parts = [
    (STEEL,  "Part 1", "Own State",
     "3D position (x, y, z) + velocity vector + current heading"),
    (GREEN,  "Part 2", "Assignment Information",
     "3D relative vector to assigned target\nOutput of Hungarian algorithm — updated every step"),
    (ORANGE, "Part 3", "Conflict Neighborhood",
     "Positions and velocities of drones predicted\nto collide with me (sparse conflict graph output)"),
    (PURPLE, "Part 4", "Obstacle Proximity",
     "Distance readings in 6 cardinal directions:\nNorth, South, East, West, Up, Down"),
]
for i, (col, num, label, desc) in enumerate(parts):
    y = 2.02 + i * 1.28
    rect(s, 0.50, y, 0.55, 1.10, fill=col)
    tb(s, num, 0.50, y, 0.55, 1.10, sz=12, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)
    rect(s, 1.10, y, 5.00, 1.10, fill=LGRAY, line=col, lw=Pt(0.8))
    tb(s, label, 1.18, y+0.04, 4.85, 0.38, sz=13, bold=True, color=col)
    tb(s, desc,  1.18, y+0.44, 4.85, 0.60, sz=12, color=DGRAY)

# Pipeline (right)
rect(s, 6.50, 1.52, 6.55, 5.75, fill=WHITE, line=STEEL, lw=Pt(1.2))
section_pill(s, "Framework Pipeline", 6.50, 1.52, 6.55, 0.40, NAVY)

steps = [
    (STEEL,                    "3D Simulation Environment (PyBullet)"),
    (GREEN,                    "Hungarian Algorithm (per step) → Assignment"),
    (ORANGE,                   "Conflict Graph (per step) → Neighbors"),
    (NAVY,                     "Combined Observation Vector (4 parts)"),
    (PURPLE,                   "MAPPO Actor Network (per drone, local obs.)"),
    (RGBColor(0xC6,0x28,0x28), "3D Velocity Action Executed"),
    (MGRAY,                    "MAPPO Centralized Critic (training only)"),
]
for i, (col, label) in enumerate(steps):
    y = 2.05 + i * 0.72
    rect(s, 6.75, y, 6.05, 0.55, fill=col)
    tb(s, label, 6.75, y, 6.05, 0.55, sz=13, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        tb(s, "↓", 9.5, y+0.55, 0.4, 0.20, sz=12, bold=True, color=NAVY,
           align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – TRAINING STRATEGY
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Training Strategy: 4-Stage Curriculum Learning",
       "Validated by both DA-MAPPO and IGAT-MARL — difficulty increases progressively")

tbl(s,
    ["Stage", "Drones", "Targets", "Obstacles", "Environment", "Purpose"],
    [
        ["1\nBaseline",       "3", "Static",  "30",       "3D — simple",        "Replicate DA-MAPPO in 3D; establish validation baseline"],
        ["2\nInteraction",    "5", "Moving",  "30-40",    "3D — moderate",      "First exposure to combined assignment-conflict interference"],
        ["3\nFull Stress",    "8", "Dynamic", "40-50",    "3D — high density",  "Maximum simultaneous interference; primary experiment"],
        ["4\nGeneralization", "3-8","Dynamic","Mixed",    "3D — unseen configs","Test transfer to swarm sizes not seen during training"],
    ],
    0.40, 1.55, W-0.8, 3.05, fs=13)

rect(s, 0.40, 4.70, W-0.8, 0.38, fill=STEEL)
tb(s, "Why Curriculum Learning?", 0.40, 4.70, W-0.8, 0.38,
   sz=14, bold=True, color=WHITE)

info_box(s, "Justification",
    ["Avoids catastrophic policy failure when exposed to hard cases too early",
     "Both DA-MAPPO and IGAT-MARL independently validated this strategy",
     "Difficulty is increased only when the policy converges at the current stage",
    ], 0.40, 5.15, 6.3, 2.15, sz=13)

info_box(s, "Simulation & Tools",
    ["Simulator: PyBullet (open-source 3D physics engine)",
     "Policy framework: PyTorch MAPPO implementation",
     "Hardware: Training on GPU; inference tested on CPU",
     "All tools open-source — fully reproducible experiments",
    ], 6.85, 5.15, 6.10, 2.15, sz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – EVALUATION PLAN
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Evaluation Plan",
       "Controlled experiments across swarm sizes, obstacle densities, and baselines")

section_pill(s, "Primary Metric", 0.40, 1.52, 3.0, 0.38, NAVY)
rect(s, 0.40, 1.95, 3.0, 1.25, fill=WHITE, line=NAVY, lw=Pt(1.5))
tb(s, "Mission Success Rate",
   0.50, 1.98, 2.85, 0.42, sz=14, bold=True, color=NAVY)
tb(s, "% of episodes where ALL drones reach their targets with zero collisions within the time limit",
   0.50, 2.42, 2.85, 0.75, sz=12, color=DGRAY)

section_pill(s, "Secondary Metrics", 3.60, 1.52, 4.5, 0.38, STEEL)
bullets(s, ["Inter-drone collision count per episode",
            "Obstacle collision count per episode",
            "Target reassignments per episode",
            "Average trajectory length per drone"],
        3.60, 1.95, 4.5, 1.25, sz=13)

section_pill(s, "Test Configurations", 8.35, 1.52, 4.65, 0.38, GREEN)
tbl(s,
    ["Swarm Sizes", "Obstacle Densities"],
    [["3, 5, 8 drones", "30, 40, 50 obstacles"]],
    8.35, 1.95, 4.65, 0.90, fs=13)
tb(s, "= 9 conditions × 4 baselines = 36 experiment runs",
   8.35, 2.92, 4.65, 0.30, sz=12, bold=True, color=GREEN)

# baselines
rect(s, 0.40, 3.35, W-0.8, 0.38, fill=NAVY)
tb(s, "Four Controlled Baselines", 0.40, 3.35, W-0.8, 0.38,
   sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

baselines = [
    (STEEL,  "B1", "Standard MAPPO",
     "MAPPO with no assignment mechanism and no conflict graph.",
     "Measures: how much does each mechanism contribute individually?"),
    (GREEN,  "B2", "DA-MAPPO ported to 3D",
     "Assignment only, no conflict graph, extended to 3D.",
     "Measures: what does the conflict graph add over assignment alone?"),
    (ORANGE, "B3", "IGAT-MARL + fixed assignment",
     "Conflict graph only, assignment is fixed at episode start.",
     "Measures: what does dynamic reassignment add over fixed assignment?"),
    (PURPLE, "B4", "Original 2D DA-MAPPO",
     "Exact replication of Sheng et al. in 2D.",
     "Measures: does this research reproduce the original results?"),
]
for i, (col, num, name, desc, purpose) in enumerate(baselines):
    x = 0.40 + i * 3.22
    rect(s, x, 3.82, 3.08, 3.52, fill=WHITE, line=col, lw=Pt(1.5))
    rect(s, x, 3.82, 3.08, 0.40, fill=col)
    tb(s, f"{num}: {name}", x+0.08, 3.82, 2.95, 0.40,
       sz=12, bold=True, color=WHITE)
    tb(s, desc,    x+0.10, 4.28, 2.90, 1.30, sz=12, color=DGRAY)
    rect(s, x+0.08, 5.65, 2.93, 0.05, fill=LGRAY)
    tb(s, purpose, x+0.10, 5.72, 2.90, 1.50, sz=11, italic=True, color=MGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – EXPECTED CONTRIBUTIONS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Expected Contributions",
       "What this research will add to the field of multi-UAV coordination")

contribs = [
    (STEEL,  "1",
     "First Combined Framework",
     "The first architecture that integrates dynamic target assignment (Hungarian algorithm) and conflict-aware collision avoidance (sparse interaction graph) into a single MARL policy for multi-UAV coordination.",
     "Fills the gap documented by both DA-MAPPO and IGAT-MARL in their own future work sections."),
    (GREEN,  "2",
     "First Empirical Evidence in 3D",
     "The first experimental results measuring whether assignment and collision-avoidance mechanisms cooperate or compete when both are encoded in the same observation vector inside a 3D environment.",
     "Answers a question that no prior work has addressed, despite 3D being the real deployment environment."),
    (ORANGE, "3",
     "Unified Observation Design",
     "A reusable observation vector design that encodes assignment state, conflict neighborhood, own state, and obstacle proximity — applicable to future multi-UAV coordination research.",
     "Provides a principled starting point for researchers extending this work to larger swarms or real hardware."),
    (PURPLE, "4",
     "Failure Boundary Characterization",
     "A systematic map of where and how the combined policy fails — across swarm sizes, obstacle densities, and target speeds — with characterization of the failure mode.",
     "This boundary is as scientifically valuable as success: it tells future researchers exactly where to focus next."),
]
for i, (col, num, title, desc, impact) in enumerate(contribs):
    y = 1.52 + i * 1.42
    rect(s, 0.40, y, 0.72, 1.22, fill=col)
    tb(s, num, 0.40, y, 0.72, 1.22, sz=24, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)
    rect(s, 1.18, y, W-1.58, 1.22, fill=WHITE, line=col, lw=Pt(1.2))
    tb(s, title, 1.30, y+0.05, W-1.80, 0.38, sz=14, bold=True, color=col)
    tb(s, desc,  1.30, y+0.43, W-1.80, 0.48, sz=12, color=DGRAY)
    tb(s, f"Impact: {impact}", 1.30, y+0.88, W-1.80, 0.35, sz=11,
       italic=True, color=MGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 – REFERENCES
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "References",
       "Papers reviewed in this research (papers present in study folder)")

refs = [
    "[1]  J. Tang, Y. Liang & K. Li, \"Dynamic scene path planning of UAVs based on deep reinforcement learning,\" Drones, vol. 8, no. 2, p. 60, 2024.",
    "[2]  X. Kong, Y. Zhou, Z. Li & S. Wang, \"Multi-UAV simultaneous target assignment and path planning based on DRL,\" Front. Neurorobot., vol. 17, p. 1302898, 2024.",
    "[3]  R. Jarray, I. Zaghbani & S. Bouallègue, \"Dynamic reward-based DRL algorithm for UAV path planning in large-scale environments,\" Procedia Comput. Sci., vol. 270, pp. 692-702, 2025.",
    "[4]  Y. Zhang et al., \"Large-scale UAV swarm path planning based on mean-field reinforcement learning,\" Chin. J. Aeronaut., vol. 38, no. 9, p. 103484, 2025.",
    "[6]  Y. Xu et al., \"Scalable UAV multi-hop networking via multi-agent reinforcement learning with large language models,\" arXiv:2505.08448, 2026.",
    "[7]  Y. Wang et al., \"RALLY: Role-adaptive LLM-driven yoked navigation for agentic UAV swarms,\" IEEE Open J. Veh. Technol., vol. 6, 2025.",
    "[8]  M. S. Khan et al., \"Beyond single-framework architectures: A systematic evaluation and hybrid design for scalable multi-agent coordination,\" IEEE Access, vol. 14, 2026.",
    "[9]  M. R. Rezaee et al., \"Efficient multi-agent DRL algorithm for multi-UAV collision avoidance,\" Appl. Soft Comput., vol. 197, p. 115145, 2026.",
    "[10] Y. Sheng, X. Xie, H. Liu & J. Li, \"Dynamic target assignment and cooperative decision-making for UAV swarms based on MARL,\" IEEE Internet Things J., 2026.",
    "[11] S. Govinda, B. Brik & S. Harous, \"A survey on deep reinforcement learning applications in autonomous systems,\" IEEE Trans. Intell. Transp. Syst., vol. 26, no. 7, 2025.",
]
left_refs  = refs[:5]
right_refs = refs[5:]

for i, ref in enumerate(left_refs):
    y = 1.60 + i * 1.06
    rect(s, 0.40, y, W/2 - 0.55, 0.95, fill=WHITE, line=STEEL, lw=Pt(0.8))
    tb(s, ref, 0.55, y+0.05, W/2 - 0.75, 0.88, sz=11, color=DGRAY)

for i, ref in enumerate(right_refs):
    y = 1.60 + i * 1.06
    rect(s, W/2 + 0.15, y, W/2 - 0.55, 0.95, fill=WHITE, line=STEEL, lw=Pt(0.8))
    tb(s, ref, W/2 + 0.30, y+0.05, W/2 - 0.75, 0.88, sz=11, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 – THANK YOU
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=DARK)
rect(s, 0, 0, W, H*0.52, fill=NAVY)
rect(s, 0, H*0.52-0.05, W, 0.10, fill=GOLD)

tb(s, "Thank You", 0.5, 0.35, W-1, 1.5,
   sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Questions & Discussion", 0.5, 1.85, W-1, 0.65,
   sz=22, color=GOLD, align=PP_ALIGN.CENTER)

rect(s, 1.5, 2.80, W-3, 0.52, fill=GOLD)
tb(s,
   "First study of assignment-avoidance interference in 3D multi-UAV coordination",
   1.5, 2.80, W-3, 0.52, sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

mtb(s, [
    ("Ayesha Khalil   |   SP25-RCS-009", True,  WHITE, 18, PP_ALIGN.CENTER),
    ("MS Computer Science — COMSATS University Islamabad", False, LIGHT_BLU, 14, PP_ALIGN.CENTER),
], 0.5, 4.15, W-1, 0.95)

# ════════════════════════════════════════════════════════════════════════════
out = "/Users/rightsteps/Masters/deep reinforcement learning/ayesha/Synopsis_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
