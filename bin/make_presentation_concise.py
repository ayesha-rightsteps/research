from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x38, 0x64)
STEEL  = RGBColor(0x2E, 0x75, 0xB6)
LBLUE  = RGBColor(0xDF, 0xE9, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xF0, 0xA5, 0x00)
DARK   = RGBColor(0x12, 0x12, 0x1E)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
MGRAY  = RGBColor(0x77, 0x77, 0x77)
LGRAY  = RGBColor(0xF5, 0xF7, 0xFC)
GREEN  = RGBColor(0x1B, 0x5E, 0x20)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)
RED    = RGBColor(0xBF, 0x36, 0x0C)
LRED   = RGBColor(0xFF, 0xF3, 0xE0)
PURP   = RGBColor(0x4A, 0x14, 0x8C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H  = 13.33, 7.5

# ── Primitives ────────────────────────────────────────────────────────────────
def R(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    sh.line.width = lw
    if line: sh.line.color.rgb = line
    else:    sh.line.fill.background()
    return sh

def T(slide, text, l, t, w, h, sz=18, bold=False, color=DGRAY,
      align=PP_ALIGN.LEFT, italic=False):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    bx.word_wrap = True
    tf = bx.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return bx

def MT(slide, lines, l, t, w, h, dsz=16, dc=DGRAY, da=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str): txt,b,c,s,a = item,False,dc,dsz,da
        else:
            txt=item[0]; b=item[1] if len(item)>1 else False
            c=item[2] if len(item)>2 else dc; s=item[3] if len(item)>3 else dsz
            a=item[4] if len(item)>4 else da
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment = a; r=p.add_run(); r.text=txt
        r.font.size=Pt(s); r.font.bold=b; r.font.color.rgb=c

def hdr(slide, title, sub=None):
    R(slide, 0, 0, W, 1.18, fill=NAVY)
    R(slide, 0, 1.18, W, 0.05, fill=GOLD)
    T(slide, title, 0.4, 0.10, W-0.8, 0.72, sz=24, bold=True, color=WHITE)
    if sub:
        T(slide, sub, 0.4, 0.78, W-0.8, 0.38, sz=13, color=GOLD)

def bg(slide): R(slide, 0, 0, W, H, fill=LGRAY)

def stat_box(slide, number, label, l, t, w=2.8, h=1.5,
             num_col=GOLD, bg_col=NAVY):
    R(slide, l, t, w, h, fill=bg_col)
    T(slide, number, l, t,   w, h*0.60, sz=34, bold=True, color=num_col,
      align=PP_ALIGN.CENTER)
    T(slide, label,  l, t+h*0.58, w, h*0.42, sz=12, color=WHITE,
      align=PP_ALIGN.CENTER)

def icon_row(slide, items, y, col=STEEL):
    """items = list of (icon_text, label)  — evenly spaced"""
    n   = len(items)
    bw  = (W - 0.8) / n
    for i, (icon, label) in enumerate(items):
        x = 0.4 + i * bw
        R(slide, x+0.1, y, bw-0.2, 1.0, fill=col)
        T(slide, icon,  x+0.1, y,       bw-0.2, 0.55, sz=22, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
        T(slide, label, x+0.1, y+0.55,  bw-0.2, 0.45, sz=12, color=WHITE,
          align=PP_ALIGN.CENTER)

def card(slide, title, points, l, t, w, h, hcol=STEEL, bcol=WHITE, sz=14):
    R(slide, l, t, w, h, fill=bcol, line=hcol, lw=Pt(1.2))
    R(slide, l, t, w, 0.38, fill=hcol)
    T(slide, title, l, t, w, 0.38, sz=13, bold=True, color=WHITE,
      align=PP_ALIGN.CENTER)
    lines = []
    for p in points:
        lines.append((f"•  {p}", False, DGRAY, sz, PP_ALIGN.LEFT))
    MT(slide, lines, l+0.12, t+0.44, w-0.24, h-0.50, dsz=sz)

def divider(slide, y, col=GOLD):
    R(slide, 0.4, y, W-0.8, 0.05, fill=col)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, W, H, fill=DARK)
R(s, 0, 0, W, H*0.56, fill=NAVY)
R(s, 0, H*0.56-0.05, W, 0.10, fill=GOLD)

T(s, "COMSATS University Islamabad",
  1, 0.18, W-2, 0.45, sz=13, color=LBLUE, align=PP_ALIGN.CENTER)
T(s, "Dynamic Target Assignment\n& Collision Avoidance\nfor Multi-UAV Coordination in 3D",
  0.8, 0.65, W-1.6, 2.4, sz=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(s, "using Multi-Agent Reinforcement Learning",
  0.8, 2.95, W-1.6, 0.5, sz=17, color=GOLD, align=PP_ALIGN.CENTER)

MT(s, [
    ("Ayesha Khalil  |  SP25-RCS-009", True,  WHITE, 18, PP_ALIGN.CENTER),
    ("MS Computer Science", False, LBLUE, 14, PP_ALIGN.CENTER),
], 0.8, 4.55, W-1.6, 1.0)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE  (visual numbered list)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Agenda")

items = [
    (STEEL, "01", "Introduction"),
    (STEEL, "02", "Motivation"),
    (STEEL, "03", "Background"),
    (GREEN, "04", "Literature Review"),
    (GREEN, "05", "Most Related Work"),
    (RED,   "06", "Problem Statement"),
    (RED,   "07", "Objectives"),
    (PURP,  "08", "Methodology"),
    (PURP,  "09", "Evaluation"),
    (NAVY,  "10", "Contributions"),
]
cols = 5; bw = (W-0.8)/cols; bh = 1.38
for i, (col, num, label) in enumerate(items):
    c = i % cols; r = i // cols
    x = 0.4 + c*bw; y = 1.35 + r*bh
    R(s, x+0.08, y, bw-0.16, bh-0.12, fill=WHITE, line=col, lw=Pt(1.5))
    R(s, x+0.08, y, 0.58,    bh-0.12, fill=col)
    T(s, num, x+0.08, y, 0.58, bh-0.12, sz=15, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, label, x+0.72, y+0.38, bw-0.88, 0.65, sz=13, bold=True, color=col)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Introduction", "Why does multi-UAV coordination matter?")

icon_row(s, [
    ("🚨", "Disaster Response"),
    ("🔍", "Search & Rescue"),
    ("🏗", "Infrastructure Inspection"),
    ("📡", "Surveillance"),
], y=1.42, col=STEEL)

divider(s, 2.62)

T(s, "The Challenge", 0.4, 2.75, W-0.8, 0.40, sz=15, bold=True, color=NAVY)

challenges = [
    (STEEL,  "Navigate",  "Reach assigned targets"),
    (RED,    "Avoid",     "No mid-air collisions"),
    (GREEN,  "Coordinate","Share airspace with the team"),
]
bw2 = (W-0.8)/3
for i, (col, title, sub) in enumerate(challenges):
    x = 0.4 + i*bw2
    R(s, x+0.1, 3.22, bw2-0.2, 1.10, fill=col)
    T(s, title, x+0.1, 3.22, bw2-0.2, 0.58, sz=20, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, sub,   x+0.1, 3.80, bw2-0.2, 0.52, sz=13,
      color=WHITE, align=PP_ALIGN.CENTER)

divider(s, 4.50)

MT(s, [
    ("Classical planners (A*, RRT): ", True, NAVY, 15),
    ("fail when environment changes mid-flight", False, DGRAY, 15),
], 0.4, 4.65, 6.3, 0.4)
MT(s, [
    ("Deep Reinforcement Learning: ", True, GREEN, 15),
    ("learns adaptive policies from simulation", False, DGRAY, 15),
], 0.4, 5.12, 6.3, 0.4)
MT(s, [
    ("Multi-Agent RL (MARL): ", True, STEEL, 15),
    ("extends DRL to cooperating drone teams", False, DGRAY, 15),
], 0.4, 5.59, 6.3, 0.4)
MT(s, [
    ("This research: ", True, RED, 15),
    ("target assignment  +  collision avoidance  +  3D  — never done together", False, DGRAY, 15),
], 0.4, 6.06, W-0.8, 0.4)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MOTIVATION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Motivation", "Why this problem? Why now?")

# Left: the gap in one sentence
R(s, 0.4, 1.38, W-0.8, 0.80, fill=NAVY)
T(s, "Two sub-problems solved independently for years — but never together, never in 3D",
  0.6, 1.42, W-1.2, 0.72, sz=18, bold=True, color=GOLD,
  align=PP_ALIGN.CENTER)

# Three motivation pillars
pillars = [
    (STEEL, "Real-World Gap",
     "UAVs fly in 3D.\nAll multi-UAV MARL\nworks in 2D.",
     "Every warehouse, rescue mission,\nand delivery fleet operates\nin three dimensions."),
    (RED,   "Documented Gap",
     "DA-MAPPO and IGAT-MARL\nboth cite each other\nas future work.",
     "The gap is not an opinion —\nit is documented by the\nauthors themselves."),
    (GREEN, "Scientific Gap",
     "No one knows if these\ntwo mechanisms cooperate\nor compete.",
     "This is an open empirical\nquestion with real consequences\nfor mission success."),
]
bwp = (W-0.8)/3
for i, (col, title, key, detail) in enumerate(pillars):
    x = 0.4 + i*bwp
    R(s, x+0.1, 2.38, bwp-0.2, 4.90, fill=WHITE, line=col, lw=Pt(2))
    R(s, x+0.1, 2.38, bwp-0.2, 0.52, fill=col)
    T(s, title, x+0.1, 2.38, bwp-0.2, 0.52, sz=15, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, key,   x+0.22, 3.00, bwp-0.44, 1.60, sz=15, bold=True, color=col)
    R(s, x+0.22, 4.65, bwp-0.44, 0.04, fill=LGRAY)
    T(s, detail, x+0.22, 4.78, bwp-0.44, 2.40, sz=13, color=MGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BACKGROUND
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Background", "Key concepts used in this research")

# DRL loop (simple 3-box)
for col, label, bx, by in [
    (NAVY,  "Agent\n(Neural Network)", 0.5,  3.0),
    (STEEL, "Environment\n(3D Simulation)", 4.8, 3.0),
]:
    R(s, bx, by, 3.2, 1.15, fill=col)
    T(s, label, bx, by, 3.2, 1.15, sz=15, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
T(s, "→ Action",  3.75, 3.05, 1.0, 0.4, sz=12, color=MGRAY)
T(s, "← State",   3.75, 3.55, 1.0, 0.4, sz=12, color=MGRAY)
T(s, "↑ Reward",  1.9,  4.22, 1.2, 0.35, sz=12, color=MGRAY)
R(s, 1.9, 4.58, 0.05, 0.05)  # dummy

# CTDE
R(s, 0.4, 1.38, W-0.8, 0.05, fill=GOLD)
card(s, "CTDE Framework  (Centralized Training, Decentralized Execution)",
    ["Training: critic sees all agents' states",
     "Execution: each drone acts on its own observation only",
     "Solves non-stationarity in multi-agent learning"],
    0.4, 1.48, W-0.8, 1.42, hcol=NAVY, sz=14)

# right side concepts
card(s, "MAPPO — policy used in this research",
    ["PPO extended to multi-agent cooperative settings",
     "Stable, reproducible, works with continuous actions",
     "Used by DA-MAPPO — proven for UAV assignment"],
    8.30, 2.90, 4.75, 1.80, hcol=GREEN, sz=13)

card(s, "Hungarian Algorithm",
    ["Optimal 1-to-1 drone-to-target assignment",
     "Runs every decision step",
     "Result goes directly into each drone's observation"],
    8.30, 4.85, 4.75, 1.65, hcol=STEEL, sz=13)

card(s, "Dec-POMDP",
    ["Each drone has partial local observation",
     "Formal framework for all MARL UAV problems"],
    8.30, 6.60, 4.75, 1.08, hcol=PURP, sz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LITERATURE REVIEW (Group 1 + 2)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Literature Review", "Papers in this study — grouped by contribution")

R(s, 0.4, 1.35, W-0.8, 0.38, fill=NAVY)
T(s, "Single-UAV Navigation", 0.4, 1.35, 6.3, 0.38,
  sz=13, bold=True, color=WHITE)
T(s, "Multi-UAV Coordination", 0.4+6.35, 1.35, 6.25, 0.38,
  sz=13, bold=True, color=WHITE)

# Paper cards row 1
paper_cards = [
    (STEEL, "[1] Tang et al., 2024",
     ["Improved D3QN + Prioritized Replay",
      "Single UAV, 2D, moving obstacles",
      "95% success rate",
      "Limit: one drone, no team"]),
    (STEEL, "[3] Jarray et al., 2025",
     ["DQN + 3D CNN, 25 km² environment",
      "Single UAV, genuine 3D",
      "98% success (low obstacles)",
      "Limit: one drone only"]),
    (GREEN, "[2] Kong et al., 2024",
     ["TANet-TD3, Hungarian supervision",
      "5 UAVs, 2D, partial observability",
      "Joint assignment + path planning",
      "Limit: 5 drones, 2D only"]),
    (GREEN, "[4] Zhang et al., 2025",
     ["Mean-field DDPG + attention",
      "Up to 120 UAVs, 2D",
      ">90% success at 120 drones",
      "Limit: 2D, static obstacles"]),
]
bw3 = (W-0.8)/4
for i, (col, title, pts) in enumerate(paper_cards):
    card(s, title, pts, 0.4+i*bw3, 1.80, bw3-0.15, 2.55, hcol=col, sz=12)

divider(s, 4.52)

R(s, 0.4, 4.60, W-0.8, 0.38, fill=NAVY)
T(s, "LLM-Enhanced Coordination", 0.4, 4.60, 8.5, 0.38, sz=13, bold=True, color=WHITE)
T(s, "Survey", 0.4+8.55, 4.60, 4.2, 0.38, sz=13, bold=True, color=WHITE)

paper_cards2 = [
    (RED, "[6] Xu et al., 2026",
     ["IPPO + GPT-4o knowledge distillation",
      "12-24 UAVs, 2D networking",
      "+52% data rate",
      "Limit: 2D, GPT-4o dependency"]),
    (RED, "[7] Wang et al., 2025",
     ["RALLY: LLM roles + RMIX mixing",
      "Commander / Coordinator / Executor",
      "Generalizes to unseen swarm sizes",
      "Limit: 14.5s inference latency"]),
    (RED, "[8] Khan et al., 2026",
     ["CrewAI vs AutoGen vs LangChain",
      "96.1% success, 76% less tokens",
      "First fair framework comparison",
      "Limit: wildfire domain only"]),
    (PURP, "[11] Govinda et al., 2025",
     ["IEEE TITS survey on DRL systems",
      "Transportation, robotics, UAVs",
      "Key gap: no unified framework",
      "Directly motivates this research"]),
]
for i, (col, title, pts) in enumerate(paper_cards2):
    card(s, title, pts, 0.4+i*bw3, 5.05, bw3-0.15, 2.32, hcol=col, sz=12)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MOST RELATED: DA-MAPPO
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Most Related Work: DA-MAPPO  [10]",
           "Sheng et al. (2026) — IEEE Internet of Things Journal")

# left: key facts
card(s, "What it does",
    ["Hungarian assignment at EVERY decision step",
     "Assignment result added to each drone's observation",
     "MAPPO policy with 4-tier reward structure",
     "Curriculum training: 3 difficulty levels"],
    0.4, 1.38, 5.8, 2.42, hcol=STEEL, sz=14)

card(s, "Tested on",
    ["3 UAVs  |  3 dynamic targets  |  30-50 obstacles",
     "2D environment only",
     "Validated on NVIDIA Jetson Orin Nano"],
    0.4, 3.95, 5.8, 1.72, hcol=MGRAY, sz=14)

card(s, "What it leaves open",
    ["No 3D extension",
     "Only 3 drones — scalability unknown",
     "No collision-avoidance mechanism",
     "Authors list 3D + larger swarms as future work"],
    0.4, 5.80, 5.8, 1.58, hcol=RED, sz=14)

# right: the key finding
R(s, 6.55, 1.38, 6.45, 5.98, fill=WHITE, line=GOLD, lw=Pt(2))
R(s, 6.55, 1.38, 6.45, 0.42, fill=GOLD)
T(s, "The Most Important Finding", 6.55, 1.38, 6.45, 0.42,
  sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

T(s, "With assignment:", 6.75, 1.95, 6.1, 0.38,
  sz=13, bold=True, color=GREEN)
R(s, 6.75, 2.38, 6.1, 0.82, fill=GREEN)
T(s, "90 – 99%  mission success", 6.75, 2.38, 6.1, 0.82,
  sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

R(s, 6.75, 3.35, 6.1, 0.05, fill=LGRAY)

T(s, "Without assignment:", 6.75, 3.52, 6.1, 0.38,
  sz=13, bold=True, color=RED)
R(s, 6.75, 3.95, 6.1, 0.82, fill=RED)
T(s, "0%  mission success", 6.75, 3.95, 6.1, 0.82,
  sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

T(s, "The assignment mechanism\nis the entire reason it works.",
  6.75, 4.98, 6.1, 0.85, sz=15, bold=True, color=DARK,
  align=PP_ALIGN.CENTER)

R(s, 6.55, 6.02, 6.45, 0.05, fill=GOLD)
T(s, "Removing it = complete failure", 6.55, 6.10, 6.45, 0.55,
  sz=14, color=RED, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MOST RELATED: IGAT-MARL
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Most Related Work: IGAT-MARL  [9]",
           "Rezaee et al. (2026) — Applied Soft Computing")

card(s, "What it does",
    ["Sparse conflict graph — only collision-risk pairs connected",
     "Improved Graph Attention Network (stacked double attention)",
     "Updated every decision step",
     "Curriculum: 3 to 10 drones"],
    0.4, 1.38, 5.8, 2.42, hcol=STEEL, sz=14)

card(s, "Tested on",
    ["3-10 UAVs  |  BlueSky simulator (fixed-wing)",
     "Conflict-guaranteed episodes",
     "Simulation only"],
    0.4, 3.95, 5.8, 1.55, hcol=MGRAY, sz=14)

card(s, "What it leaves open",
    ["No target assignment — drones have no goal",
     "Authors explicitly list task allocation as future direction"],
    0.4, 5.65, 5.8, 1.72, hcol=RED, sz=14)

# right: results
R(s, 6.55, 1.38, 6.45, 5.98, fill=WHITE, line=GREEN, lw=Pt(2))
R(s, 6.55, 1.38, 6.45, 0.42, fill=GREEN)
T(s, "Results vs. Best Baseline", 6.55, 1.38, 6.45, 0.42,
  sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

stats = [("17%", "higher reward"), ("10%", "fewer dangerous events"), ("44%", "fewer graph edges")]
for i, (num, label) in enumerate(stats):
    y = 2.00 + i*1.40
    R(s, 6.75, y, 6.1, 1.15, fill=GREEN)
    T(s, num,   6.75, y,      6.1, 0.68, sz=36, bold=True,
      color=GOLD, align=PP_ALIGN.CENTER)
    T(s, label, 6.75, y+0.68, 6.1, 0.47, sz=14,
      color=WHITE, align=PP_ALIGN.CENTER)

T(s, "44% fewer edges = less communication overhead\n= scales to larger swarms",
  6.75, 6.12, 6.1, 0.72, sz=13, italic=True, color=MGRAY,
  align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Problem Statement", "What happens when you combine both in 3D?")

# gap table — compact
from pptx.util import Inches as In
cols_h = ["", "DA-MAPPO", "IGAT-MARL", "This Research"]
rows_d = [
    ["Target Assignment",   "✔",  "✘",  "✔"],
    ["Collision Avoidance", "✘",  "✔",  "✔"],
    ["3D Environment",      "✘",  "✘",  "✔"],
    ["Swarm Size",          "3",  "3-10","3,5,8"],
]
tbl_sh = s.shapes.add_table(len(rows_d)+1, len(cols_h),
    In(0.4), In(1.38), In(W-0.8), In(2.38))
tbl = tbl_sh.table
col_widths = [In(2.2), In(3.1), In(3.1), In(3.4)]
for c in range(len(cols_h)):
    tbl.columns[c].width = col_widths[c]
def cfmt(cell, text, bg_col, fg_col, bold=False, fs=13):
    cell.text = ""
    cell.fill.solid(); cell.fill.fore_color.rgb = bg_col
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = fg_col
for c, h in enumerate(cols_h):
    cfmt(tbl.cell(0,c), h, NAVY, WHITE, bold=True)
row_fills = [WHITE, LGRAY, WHITE, LGRAY]
for r, row in enumerate(rows_d):
    for c, val in enumerate(row):
        if c == 3:
            col_v = GREEN if val in ("✔","3,5,8") else RED
            cfmt(tbl.cell(r+1,c), val, LGREEN, GREEN, bold=True)
        elif val == "✘":
            cfmt(tbl.cell(r+1,c), val, LRED, RED)
        elif val == "✔":
            cfmt(tbl.cell(r+1,c), val, LGREEN, GREEN, bold=True)
        else:
            cfmt(tbl.cell(r+1,c), val, row_fills[r], DGRAY)

# interference boxes
R(s, 0.4, 3.90, W-0.8, 0.38, fill=NAVY)
T(s, "The Interference — Two modules, one policy, no shared awareness",
  0.4, 3.90, W-0.8, 0.38, sz=14, bold=True, color=GOLD)

R(s, 0.4, 4.35, 5.9, 2.95, fill=LRED, line=RED, lw=Pt(1.5))
R(s, 0.4, 4.35, 5.9, 0.40, fill=RED)
T(s, "Assignment Module", 0.4, 4.35, 5.9, 0.40, sz=13, bold=True,
  color=WHITE, align=PP_ALIGN.CENTER)
T(s, '"Go to your target now"', 0.55, 4.85, 5.6, 0.42,
  sz=15, italic=True, bold=True, color=RED)
T(s, "Unaware of which other drone pairs are in conflict",
  0.55, 5.32, 5.6, 0.42, sz=13, color=DGRAY)
R(s, 0.55, 5.85, 5.6, 0.05, fill=MGRAY)
T(s, "In 3D: drone ascending to target above it\nmay collide with descending drone — unseen",
  0.55, 5.95, 5.6, 1.20, sz=13, bold=True, color=RED)

R(s, 7.0, 4.35, 5.95, 2.95, fill=LBLUE, line=STEEL, lw=Pt(1.5))
R(s, 7.0, 4.35, 5.95, 0.40, fill=STEEL)
T(s, "Conflict Graph", 7.0, 4.35, 5.95, 0.40, sz=13, bold=True,
  color=WHITE, align=PP_ALIGN.CENTER)
T(s, '"Change course — collision predicted"', 7.15, 4.85, 5.65, 0.42,
  sz=15, italic=True, bold=True, color=STEEL)
T(s, "Unaware of current target assignment",
  7.15, 5.32, 5.65, 0.42, sz=13, color=DGRAY)
R(s, 7.15, 5.85, 5.65, 0.05, fill=MGRAY)
T(s, "Reroutes drone away from conflict —\npulls it away from its assigned target",
  7.15, 5.95, 5.65, 1.20, sz=13, bold=True, color=STEEL)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESEARCH QUESTION & OBJECTIVES
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Research Question & Objectives")

R(s, 0.4, 1.38, W-0.8, 1.55, fill=NAVY)
T(s,
  "Will a single MAPPO policy with conflict-aware graph + assignment-augmented observation\nsuccessfully navigate 5-8 UAVs in 3D with >85% mission success —\nand do the two mechanisms help or hinder each other?",
  0.6, 1.42, W-1.2, 1.48, sz=16, bold=True, color=GOLD,
  align=PP_ALIGN.CENTER)

objs = [
    (STEEL,  "Obj 1", "Design combined observation vector",
             "Hungarian assignment + conflict neighborhood + own state"),
    (GREEN,  "Obj 2", "Test with 3, 5, and 8 drones",
             "Measure how performance scales with interference"),
    (RED,    "Obj 3", "Ablation experiments",
             "Isolate contribution of each mechanism individually"),
    (PURP,   "Obj 4", "Find the failure boundary",
             "Map where and how the combined policy breaks"),
]
bw4 = (W-0.8)/4
for i, (col, num, title, sub) in enumerate(objs):
    x = 0.4 + i*bw4
    R(s, x+0.08, 3.12, bw4-0.16, 3.90, fill=WHITE, line=col, lw=Pt(1.5))
    R(s, x+0.08, 3.12, bw4-0.16, 0.52, fill=col)
    T(s, num,   x+0.08, 3.12, bw4-0.16, 0.52, sz=14, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, title, x+0.18, 3.74, bw4-0.36, 0.88, sz=14, bold=True, color=col)
    T(s, sub,   x+0.18, 4.62, bw4-0.36, 2.30, sz=13, color=MGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — METHODOLOGY
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Proposed Methodology",
           "Combined observation vector + MAPPO policy in 3D")

# Observation vector — 4 horizontal parts
R(s, 0.4, 1.38, W-0.8, 0.38, fill=NAVY)
T(s, "Each Drone's Observation Vector (4 parts)",
  0.4, 1.38, W-0.8, 0.38, sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

obs_parts = [
    (STEEL, "1", "Own State",          "3D position\nvelocity\nheading"),
    (GREEN, "2", "Assignment",         "Target relative\nposition\n(Hungarian, per step)"),
    (RED,   "3", "Conflict Neighbors", "Positions of drones\npredicted to collide\nwith me"),
    (PURP,  "4", "Obstacles",          "Proximity in\n6 directions\n(N/S/E/W/Up/Down)"),
]
bw5 = (W-0.8)/4
for i, (col, num, label, desc) in enumerate(obs_parts):
    x = 0.4 + i*bw5
    R(s, x+0.08, 1.85, bw5-0.16, 2.25, fill=WHITE, line=col, lw=Pt(1.5))
    R(s, x+0.08, 1.85, bw5-0.16, 0.50, fill=col)
    T(s, f"{num}.  {label}", x+0.08, 1.85, bw5-0.16, 0.50,
      sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s, desc, x+0.18, 2.42, bw5-0.36, 1.60, sz=13, color=DGRAY)

T(s, "→", 0.4+1*bw5-0.02, 2.68, 0.28, 0.45, sz=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
T(s, "+", 0.4+2*bw5-0.02, 2.68, 0.28, 0.45, sz=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
T(s, "+", 0.4+3*bw5-0.02, 2.68, 0.28, 0.45, sz=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

divider(s, 4.28)

# Pipeline — horizontal flow
steps = [
    (STEEL, "Simulation\n(PyBullet 3D)"),
    (GREEN, "Hungarian\nAssignment"),
    (RED,   "Conflict\nGraph"),
    (NAVY,  "Combined\nObservation"),
    (PURP,  "MAPPO\nActor"),
    (MGRAY, "3D Action\nExecuted"),
]
bw6 = (W-0.8)/len(steps)
for i, (col, label) in enumerate(steps):
    x = 0.4 + i*bw6
    R(s, x+0.08, 4.42, bw6-0.16, 1.15, fill=col)
    T(s, label, x+0.08, 4.42, bw6-0.16, 1.15, sz=13, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        T(s, "→", x+bw6-0.18, 4.75, 0.35, 0.5, sz=16, bold=True,
          color=NAVY, align=PP_ALIGN.CENTER)

R(s, 0.4, 5.72, W-0.8, 0.05, fill=GOLD)

# Training stages — compact
stages = [
    (STEEL, "Stage 1", "3 drones\nStatic targets"),
    (GREEN, "Stage 2", "5 drones\nMoving targets"),
    (RED,   "Stage 3", "8 drones\nDynamic + dense"),
    (PURP,  "Stage 4", "Unseen sizes\nGeneralization"),
]
for i, (col, stage, detail) in enumerate(stages):
    x = 0.4 + i * bw5
    R(s, x+0.08, 5.85, bw5-0.16, 1.48, fill=WHITE, line=col, lw=Pt(1))
    R(s, x+0.08, 5.85, bw5-0.16, 0.45, fill=col)
    T(s, stage,  x+0.08, 5.85, bw5-0.16, 0.45, sz=13, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, detail, x+0.18, 6.34, bw5-0.36, 0.95, sz=12, color=DGRAY)
T(s, "4-Stage Curriculum Learning", 0.4, 5.65, W-0.8, 0.22,
  sz=11, bold=True, color=MGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EVALUATION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Evaluation Plan",
           "4 baselines  ×  3 swarm sizes  ×  3 obstacle densities")

# Stat boxes
stat_box(s, "4",  "Baselines",        0.40, 1.38, 2.50, 1.45)
stat_box(s, "3",  "Swarm Sizes\n(3, 5, 8 drones)", 3.10, 1.38, 2.50, 1.45)
stat_box(s, "3",  "Obstacle Levels\n(30/40/50)",    5.80, 1.38, 2.50, 1.45)
stat_box(s, "36", "Total Conditions", 8.50, 1.38, 2.50, 1.45)

stat_box(s, "✔", "Mission success\n(primary metric)",
         11.20, 1.38, 1.75, 1.45, num_col=GOLD, bg_col=GREEN)

# Baselines
R(s, 0.4, 3.00, W-0.8, 0.38, fill=STEEL)
T(s, "Four Controlled Baselines", 0.4, 3.00, W-0.8, 0.38,
  sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

baselines = [
    (STEEL, "B1", "Standard MAPPO",        "No assignment, no conflict graph"),
    (GREEN, "B2", "DA-MAPPO → 3D",         "Assignment only, ported to 3D"),
    (RED,   "B3", "IGAT-MARL + fixed asgn","Conflict graph, fixed assignment"),
    (PURP,  "B4", "Original 2D DA-MAPPO",  "Replication verification"),
]
bw7 = (W-0.8)/4
for i, (col, num, name, desc) in enumerate(baselines):
    x = 0.4 + i*bw7
    R(s, x+0.08, 3.46, bw7-0.16, 1.72, fill=WHITE, line=col, lw=Pt(1.5))
    R(s, x+0.08, 3.46, 0.55,     1.72, fill=col)
    T(s, num, x+0.08, 3.46, 0.55, 1.72, sz=16, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, name, x+0.72, 3.52, bw7-0.95, 0.55, sz=13, bold=True, color=col)
    T(s, desc, x+0.72, 4.05, bw7-0.95, 1.05, sz=12, color=MGRAY)

# Secondary metrics
R(s, 0.4, 5.30, W-0.8, 0.38, fill=NAVY)
T(s, "Secondary Metrics", 0.4, 5.30, W-0.8, 0.38, sz=14, bold=True, color=WHITE)
metrics = ["Inter-drone collisions / episode",
           "Obstacle collisions / episode",
           "Target reassignments / episode",
           "Average trajectory length"]
bwm = (W-0.8)/4
for i, m in enumerate(metrics):
    x = 0.4 + i*bwm
    R(s, x+0.08, 5.76, bwm-0.16, 0.90, fill=LBLUE)
    T(s, m, x+0.08, 5.76, bwm-0.16, 0.90, sz=13,
      color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONTRIBUTIONS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); hdr(s, "Expected Contributions")

contribs = [
    (STEEL, "1", "First Combined Framework",
     "Assignment + collision avoidance\nin a single 3D MARL policy"),
    (GREEN, "2", "First 3D Empirical Evidence",
     "Do the two mechanisms help\nor hinder each other in 3D?"),
    (RED,   "3", "Unified Observation Design",
     "Reusable vector encoding both\nmechanisms for future research"),
    (PURP,  "4", "Failure Boundary Map",
     "Where and how the combined\npolicy fails — and why"),
]
bwc = (W-0.8)/4
for i, (col, num, title, desc) in enumerate(contribs):
    x = 0.4 + i*bwc
    R(s, x+0.08, 1.38, bwc-0.16, 5.78, fill=WHITE, line=col, lw=Pt(2))
    R(s, x+0.08, 1.38, bwc-0.16, 0.82, fill=col)
    T(s, num, x+0.08, 1.38, bwc-0.16, 0.82, sz=32, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, title, x+0.18, 2.35, bwc-0.36, 1.10, sz=15, bold=True, color=col)
    R(s, x+0.18, 3.50, bwc-0.36, 0.05, fill=LGRAY)
    T(s, desc,  x+0.18, 3.65, bwc-0.36, 3.40, sz=14, color=MGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — THANK YOU
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, W, H, fill=DARK)
R(s, 0, 0, W, H*0.55, fill=NAVY)
R(s, 0, H*0.55-0.05, W, 0.10, fill=GOLD)

T(s, "Thank You", 0.5, 0.32, W-1, 1.5, sz=54, bold=True,
  color=WHITE, align=PP_ALIGN.CENTER)
T(s, "Questions & Discussion", 0.5, 1.82, W-1, 0.6, sz=22,
  color=GOLD, align=PP_ALIGN.CENTER)

R(s, 1.8, 2.75, W-3.6, 0.50, fill=GOLD)
T(s, "First study of assignment-avoidance interference in 3D multi-UAV coordination",
  1.8, 2.75, W-3.6, 0.50, sz=13, bold=True, color=DARK,
  align=PP_ALIGN.CENTER)

MT(s, [
    ("Ayesha Khalil  |  SP25-RCS-009", True,  WHITE, 17, PP_ALIGN.CENTER),
    ("MS Computer Science — COMSATS University Islamabad", False, LBLUE, 13, PP_ALIGN.CENTER),
], 0.5, 4.10, W-1, 0.90)

# ════════════════════════════════════════════════════════════════════════════
out = "/Users/rightsteps/Masters/deep reinforcement learning/ayesha/Synopsis_Presentation_Concise.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
